import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Colors matching the cyber-themed CSS tokens of the portal
    BG_COLOR = RGBColor(7, 7, 13)       # Dark Space (#07070d)
    PRIMARY_TXT = RGBColor(240, 240, 245) # Cool off-white
    CYAN_ACCENT = RGBColor(0, 240, 255)  # Electric Cyber Cyan (#00f0ff)
    GREEN_SUCCESS = RGBColor(0, 255, 102) # Matrix Green (#00ff66)
    MUTED_TXT = RGBColor(140, 145, 160)   # Muted Gray
    CARD_BG = RGBColor(15, 15, 28)       # Slightly lighter dark for card containers
    BORDER_COLOR = RGBColor(30, 41, 59)   # Slate gray border for cards

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide layout index 6 is typically a blank slide
    blank_layout = prs.slide_layouts[6]

    # Helper to set solid background color
    def set_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw accent borders
    def add_decorations(slide):
        # Thin cyan top banner accent
        top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = CYAN_ACCENT
        top_accent.line.fill.background()
        
        # Cyber grid corner decoration (aesthetic detail)
        accent_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), Inches(0.2), Inches(0.5), Inches(0.04))
        accent_box.fill.solid()
        accent_box.fill.fore_color.rgb = CYAN_ACCENT
        accent_box.line.fill.background()

    # Helper to add standard slide titles
    def add_slide_header(slide, title, section_num=""):
        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(10), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        if section_num:
            run_num = p.add_run()
            run_num.text = f"{section_num}  |  "
            run_num.font.name = "Segoe UI"
            run_num.font.size = Pt(28)
            run_num.font.bold = True
            run_num.font.color.rgb = CYAN_ACCENT
            
        run_title = p.add_run()
        run_title.text = title
        run_title.font.name = "Segoe UI"
        run_title.font.size = Pt(28)
        run_title.font.bold = True
        run_title.font.color.rgb = PRIMARY_TXT
        
        # Horizontal rule divider
        hr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(11.83), Inches(0.02))
        hr.fill.solid()
        hr.fill.fore_color.rgb = RGBColor(30, 41, 59)
        hr.line.fill.background()

    # Helper to parse and insert bullet points with formatting
    def parse_and_add_bullet(tf, text, font_size=15, line_spacing=10):
        if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.space_after = Pt(line_spacing)
        
        # Parse ** bold tags
        parts = text.split("**")
        is_bold = False
        
        # Add custom colored bullet prefix
        run_bullet = p.add_run()
        run_bullet.text = "▪  "
        run_bullet.font.color.rgb = CYAN_ACCENT
        run_bullet.font.bold = True
        
        for part in parts:
            if part == "":
                is_bold = not is_bold
                continue
            run = p.add_run()
            run.text = part
            if is_bold:
                run.font.bold = True
                run.font.color.rgb = CYAN_ACCENT
            else:
                run.font.bold = False
                run.font.color.rgb = PRIMARY_TXT
            is_bold = not is_bold

    # Helper for adding cards in grid
    def add_card(slide, left, top, width, height, title, body_bullets, accent_color=CYAN_ACCENT):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Add subtle top border accent on card
        top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = accent_color
        top_accent.line.fill.background()

        # Text container inside card
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI Semibold"
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color
        p_title.space_after = Pt(8)
        
        for bullet in body_bullets:
            p = tf.add_paragraph()
            p.font.name = "Segoe UI"
            p.font.size = Pt(12)
            p.font.color.rgb = PRIMARY_TXT
            p.space_after = Pt(6)
            
            # Simple bold formatting for card bullets
            parts = bullet.split("**")
            is_bold = False
            for part in parts:
                run = p.add_run()
                run.text = part
                if is_bold:
                    run.font.bold = True
                    run.font.color.rgb = CYAN_ACCENT
                is_bold = not is_bold

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Cyber theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    
    # Large Title text box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CYBER CRIME PORTAL"
    p.font.name = "Segoe UI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "A Secure, Role-Based, Real-Time Incident Reporting & Case Management System"
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = PRIMARY_TXT
    p_sub.alignment = PP_ALIGN.LEFT
    p_sub.space_after = Pt(20)

    # Decorative separator bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(4.3), Inches(4.5), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_SUCCESS
    bar.line.fill.background()

    # Meta text box at bottom
    meta_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.33), Inches(1.5))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_tech = tf_meta.paragraphs[0]
    p_tech.text = "Tech Stack: React 19 (Vite)  |  Node.js (Express)  |  Sequelize ORM  |  MySQL  |  Socket.IO"
    p_tech.font.name = "Segoe UI Semibold"
    p_tech.font.size = Pt(14)
    p_tech.font.color.rgb = MUTED_TXT
    
    p_desc = tf_meta.add_paragraph()
    p_desc.text = "Designed for seamless anonymous reporting and structured investigation management."
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = MUTED_TXT
    p_desc.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 2: Executive Summary
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "EXECUTIVE SUMMARY", "01")

    # Left Column (Overview text)
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_left = [
        "**Secure Reporting Platform**: Bridging the communication gap between citizens (who can file complaints officially or anonymously) and law enforcement agencies.",
        "**Real-Time Subsystem**: The application leverages WebSockets via **Socket.IO** to handle instantaneous notifications and case updates across all client instances.",
        "**Decoupled Architecture**: Fully separated client-server design allowing modular scaling, strict access controls, and robust security policies."
    ]
    for bullet in bullets_left:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=14)

    # Right Column (Key Features Callout Box)
    card_left = Inches(7.0)
    card_top = Inches(1.6)
    card_width = Inches(5.58)
    card_height = Inches(4.8)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(1.5)
    
    right_box = slide.shapes.add_textbox(card_left + Inches(0.3), card_top + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_right_title = tf_right.paragraphs[0]
    p_right_title.text = "CORE VALUE PROPOSITIONS"
    p_right_title.font.name = "Segoe UI"
    p_right_title.font.size = Pt(16)
    p_right_title.font.bold = True
    p_right_title.font.color.rgb = GREEN_SUCCESS
    p_right_title.space_after = Pt(14)
    
    bullets_right = [
        "**Privacy First**: Secure anonymous filing mode allows reporting sensitive cyber crimes without disclosing profile metadata.",
        "**Structured Cases**: Full digital chain of custody from filing to review, investigation, and final resolution.",
        "**Audit-Ready Files**: Both citizens and assigned officers can upload evidence files securely attached to individual complaints.",
        "**Interactive Dashboard**: Real-time analytical graphs, badges, and toaster alerts maintain live feedback loops."
    ]
    for bullet in bullets_right:
        parse_and_add_bullet(tf_right, bullet, font_size=13, line_spacing=8)

    # ----------------------------------------------------
    # SLIDE 3: System Architecture & Tech Stack
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "SYSTEM ARCHITECTURE & TECH STACK", "02")
    
    # Subtext explaining decoupled architecture
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.5))
    tf_desc = desc_box.text_frame
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "The platform is designed with a decoupled Client-Server architecture using modern web technologies:"
    p_desc.font.name = "Segoe UI"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = MUTED_TXT

    # Add 3 cards for Frontend, Backend, and Database
    card_width = Inches(3.64)
    card_height = Inches(4.6)
    gap = Inches(0.45)
    
    frontend_bullets = [
        "**React 19 & Vite**: Ultra-fast build and render cycles.",
        "**Context API**: Manage authentication state & active sessions.",
        "**Axios Client**: Custom interceptors handling JWT headers and error catchers.",
        "**Socket.IO Client**: Establishes persistent real-time socket connections.",
        "**Themeable styling**: Variable-based glassmorphic design system."
    ]
    add_card(slide, Inches(0.75), Inches(1.9), card_width, card_height, "FRONTEND CLIENT", frontend_bullets, CYAN_ACCENT)

    backend_bullets = [
        "**Node.js & Express**: Handles asynchronous HTTP routes & static evidence assets.",
        "**Socket.IO Server**: Manages WebSocket handshakes, state tracking, and rooms.",
        "**JWT & BCrypt**: Issues access tokens & implements secure password hashing.",
        "**Sequelize ORM**: Manages schema sync, transactions, and seeding logic."
    ]
    add_card(slide, Inches(0.75) + card_width + gap, Inches(1.9), card_width, card_height, "BACKEND API", backend_bullets, GREEN_SUCCESS)

    database_bullets = [
        "**MySQL 8.4 Server**: Reliable relational data storage engine.",
        "**Entity Relations**: Configured foreign key dependencies and cascades.",
        "**Auto Seeding**: Auto-checks database state and seeds admin, user, and officers.",
        "**Isolated Data**: Localized data folders (`mysql_data`) for portable setups."
    ]
    add_card(slide, Inches(0.75) + 2*(card_width + gap), Inches(1.9), card_width, card_height, "DATABASE ENGINE", database_bullets, CYAN_ACCENT)


    # ----------------------------------------------------
    # SLIDE 4: Database Schema & Relational Model
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "DATABASE ARCHITECTURE & SCHEMAS", "03")

    # Left TextBox for summary
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_db = [
        "**Users Table (`users`)**: Stores credentials, profile info, and accounts roles (`ROLE_CITIZEN`, `ROLE_OFFICER`, `ROLE_ADMIN`).",
        "**Complaints Table (`complaints`)**: Stores case titles, categories, priorities (`LOW`/`MEDIUM`/`HIGH`), status state, and foreign keys for author & officer.",
        "**Evidence Files (`evidence_files`)**: Tracks paths and meta of uploads attached by citizens during case creation."
    ]
    for bullet in bullets_db:
        parse_and_add_bullet(tf_left, bullet, font_size=14, line_spacing=12)

    # Right TextBox
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    bullets_db_2 = [
        "**Officer Notes Table (`officer_notes`)**: Stores detailed investigation logs added by officers. Tied to officer and case.",
        "**Evidence Uploads (`evidence_uploads`)**: Tracks documents and media uploaded by officers during the investigation phase.",
        "**Notifications Table (`notifications`)**: Stores persistent notifications (`INFO`, `WARNING`, `SUCCESS`, `ERROR`) sent to users, tracking read states."
    ]
    for bullet in bullets_db_2:
        parse_and_add_bullet(tf_right, bullet, font_size=14, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 5: Real-Time Notification Gateway (Socket.IO)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "REAL-TIME NOTIFICATION SUBSYSTEM", "04")

    # Details
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_sockets = [
        "**Direct Room Strategy**: Upon connection and JWT validation, a socket joins a room named `user_<userId>`. This allows directed alerts to a specific citizen or officer.",
        "**Role-Based Rooms**: Users also join rooms representing their roles (`admin`, `officer`, `citizen`), enabling group broadcasts.",
        "**Toast alerts**: Incoming WebSocket notifications instantly trigger animations and toast popups in the React application using `React Hot Toast`."
    ]
    for bullet in bullets_sockets:
        parse_and_add_bullet(tf_left, bullet, font_size=14, line_spacing=12)

    # Right Callout Card (WebSocket Workflow)
    card_left = Inches(6.8)
    card_top = Inches(1.6)
    card_width = Inches(5.78)
    card_height = Inches(4.8)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(1.5)
    
    inner_box = slide.shapes.add_textbox(card_left + Inches(0.3), card_top + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    tf_inner = inner_box.text_frame
    tf_inner.word_wrap = True
    
    p_inner_title = tf_inner.paragraphs[0]
    p_inner_title.text = "EVENT TRIGGERS & PAYLOAD FLOW"
    p_inner_title.font.name = "Segoe UI"
    p_inner_title.font.size = Pt(16)
    p_inner_title.font.bold = True
    p_inner_title.font.color.rgb = CYAN_ACCENT
    p_inner_title.space_after = Pt(14)
    
    triggers = [
        "**1. Complaint Submitted**: Broadcasts payload to the `admin` room. All admin badges increment instantly.",
        "**2. Case Assigned**: Direct server message sent to `user_<officerId>`. Officer dashboard updates.",
        "**3. Status Updated**: Server notifies `user_<citizenId>`. Citizen receives green/amber progress notifications.",
        "**4. Case Resolved**: Notifies both citizen and admins. Success toaster fired client-side."
    ]
    for bullet in triggers:
        parse_and_add_bullet(tf_inner, bullet, font_size=12, line_spacing=6)


    # ----------------------------------------------------
    # SLIDE 6: Citizen Portal Walkthrough
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "CITIZEN PORTAL FEATURES", "05")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_citizen = [
        "**Anonymous Reporting Mode**: Citizens can report crimes anonymously. In this mode, their identity metadata is masked in the public case view, encouraging safe whistleblowing.",
        "**Official Complaint Filing**: Comprehensive forms that collect crime categories, incident dates, descriptions, priority levels, location metadata, and file uploads.",
        "**Evidence Attachments**: Support for uploading evidence files (images, PDFs, documents) during the initial case submission.",
        "**Dynamic Interactive Dashboard**: Citizens can trace their case history, view active case statuses, and review historic notifications dynamically.",
        "**Real-Time Alerts**: Sound cues, shaking navigation bell badges, and visual toasts notify citizens instantly of officer assignments or updates."
    ]
    for bullet in bullets_citizen:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 7: Officer Portal & Investigations
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "OFFICER PORTAL & INVESTIGATIVE FLOW", "06")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_officer = [
        "**Case Queue**: A structured workflow view showing all active cyber crime complaints assigned to the logged-in officer.",
        "**Timeline Journal Notes**: Internal database logs where officers can append notes, progress details, interview summaries, and investigative findings.",
        "**Evidence Uploads**: Allows officers to upload new evidence documents, forensics files, or reports found during active investigations.",
        "**Case Status Management**: Controls to transition complaints between states: `UNDER_REVIEW` -> `INVESTIGATING` -> `RESOLVED` / `REJECTED`.",
        "**Automatic Notification Sync**: Updating status states instantly triggers websocket pushes to alert the reporting citizen."
    ]
    for bullet in bullets_officer:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 8: Administrator Portal & Dashboard
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "ADMIN PANEL & SYSTEM MANAGEMENT", "07")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_admin = [
        "**System Analytics Dashboard**: Real-time graphs showing breakdown of complaints by status (submitted, resolving, rejected), category, and priority (high, medium, low).",
        "**Complaint Assignment Panel**: Enables administrators to review new filings and assign appropriate investigation officers to cases.",
        "**User Account Controls**: Allows administrators to view citizen profiles, email logs, and toggle users' status between `ACTIVE` and `DISABLED`.",
        "**Officer Onboarding**: Secure forms to register new officers, assign default credentials, and configure investigation access clearance.",
        "**Data Audit Controls**: System rights to remove erroneous logs, manage core database syncs, and clear global notification cycles."
    ]
    for bullet in bullets_admin:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 9: UI/UX Design & Theme Styling
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "UI/UX DESIGN SYSTEM", "08")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_ui = [
        "**Cyberpunk Design Palette**: Built around modern theme variables: Space Navy Background (`#07070d`), Electric Cyan Accents (`#00f0ff`), and glowing Alert Pink (`#ff0055`).",
        "**Glassmorphism Effects**: CSS variables with background blur (`backdrop-filter: blur(12px)`) and semi-transparent border frames, creating a sleek cyber-security command center feel.",
        "**Micro-Animations**: Framer Motion handles slide transitions, list deletions, and drawer panels. Custom CSS `@keyframes shake` adds haptic notifications to the bell button.",
        "**Responsive Layouts**: Flexible grid configurations and media queries ensure the portal displays optimally across laptops, tablets, and mobile screens.",
        "**Clear Feedback Loops**: Live indicator light rings (green for active, pulse for alerts) keep users visually connected to connection states."
    ]
    for bullet in bullets_ui:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 10: Security, Authentication & Integrity
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "SECURITY & DATA PROTECTION", "09")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_security = [
        "**Password Hashing (BCryptJS)**: Solid protection using bcryptjs hashing with salt. Plain-text credentials are never written to the database or output logs.",
        "**JSON Web Token Auth**: Stateless user sessions managed via securely signed JWTs passed through request header interceptors.",
        "**Role-Based Access Control (RBAC)**: Middleware validation filters API endpoints. React route guards isolate dashboards by role (`ROLE_CITIZEN`, `ROLE_OFFICER`, `ROLE_ADMIN`).",
        "**Sequelize Parameterization**: Enforces strict database parameter mapping to prevent SQL Injection vulnerabilities.",
        "**CORS & File Validation**: Confirms evidence file uploads conform to sizes and formats, protecting backend directories from scripts execution."
    ]
    for bullet in bullets_security:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)


    # ----------------------------------------------------
    # SLIDE 11: Setup, Seeding & Execution
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "SETUP & RUNNING INSTRUCTIONS", "10")

    # Left Column (Overview text)
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_setup = [
        "**Automated Startup Script**: Windows users can start the entire project using `start-services.bat` in the root folder.",
        "**Batch Script Functions**: Installs frontend & backend nodes, initialises local `mysql_data` workspace, provisions database, seeds test credentials, and launches frontend in standard browser.",
        "**Manual Starting Commands**: If running manually:\n  - Database: Use custom port `3306` & datadir.\n  - Backend API: `npm run dev` (runs on port `8080`).\n  - React Client: `npm run dev` (runs on port `5173`)."
    ]
    for bullet in bullets_setup:
        parse_and_add_bullet(tf_left, bullet, font_size=13, line_spacing=12)

    # Right Column (Demo Credentials card)
    card_left = Inches(6.8)
    card_top = Inches(1.6)
    card_width = Inches(5.78)
    card_height = Inches(4.8)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(1.5)
    
    inner_box = slide.shapes.add_textbox(card_left + Inches(0.3), card_top + Inches(0.3), card_width - Inches(0.6), card_height - Inches(0.6))
    tf_inner = inner_box.text_frame
    tf_inner.word_wrap = True
    
    p_inner_title = tf_inner.paragraphs[0]
    p_inner_title.text = "DEMO ACCOUNTS (AUTO-SEEDED)"
    p_inner_title.font.name = "Segoe UI"
    p_inner_title.font.size = Pt(16)
    p_inner_title.font.bold = True
    p_inner_title.font.color.rgb = GREEN_SUCCESS
    p_inner_title.space_after = Pt(14)
    
    credentials = [
        "**Administrator Account**:\n  - Email: `admin@gmail.com`\n  - Password: `admin123`",
        "**Citizen Account**:\n  - Email: `user@gmail.com`\n  - Password: `user1234`",
        "**Officer Account**:\n  - Email: `officer@gmail.com`\n  - Password: `officer123`"
    ]
    for bullet in credentials:
        parse_and_add_bullet(tf_inner, bullet, font_size=12, line_spacing=8)


    # ----------------------------------------------------
    # SLIDE 12: Conclusion & Future Roadmap
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, BG_COLOR)
    add_decorations(slide)
    add_slide_header(slide, "CONCLUSION & ROADMAP", "11")

    # Content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    bullets_conclusion = [
        "**Technical Achievements**: Delivered a fully functional, real-time cyber security portal with structured workflows and secure separation of user capabilities.",
        "**High Scalability**: The decoupled client-server structure allows deploying the backend to serverless instances and database to dedicated instances without codebase changes.",
        "**Future Enhancement 1 - AI Classifiers**: Implement NLP models on the backend to automatically categorize crime reports and flag high-priority complaints based on description content.",
        "**Future Enhancement 2 - Blockchain Auditing**: Use distributed ledger technology to store file hashes for evidence submissions, ensuring tamper-proof chain of custody.",
        "**Future Enhancement 3 - In-App Chat Rooms**: Real-time communication bridge between the citizen and the assigned investigator for instant feedback."
    ]
    for bullet in bullets_conclusion:
        parse_and_add_bullet(tf_left, bullet, font_size=15, line_spacing=12)

    prs.save("Cyber_Crime_Portal_Presentation.pptx")
    print("Presentation created successfully as 'Cyber_Crime_Portal_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
